import os
import json
from pathlib import Path
from typing import List, Any, Callable, Optional, Dict
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed

try:
    import pymupdf as fitz
except ImportError:
    try:
        import fitz
    except ImportError:
        fitz = None

import numpy as np
import pandas as pd
from langchain_core.documents import Document

from config import (
    EMBEDDING_MODEL_NAME,
    EMBEDDING_DEVICE,
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    TOP_K,
    RETRIEVAL_CANDIDATES,
    OLLAMA_MODEL,
    TEMPERATURE,
)

BASE_USER_DATA_DIR = "./user_data"
_SHARED_OCR = None

def get_ocr():
    global _SHARED_OCR
    if _SHARED_OCR is None:
        from rapidocr_onnxruntime import RapidOCR
        _SHARED_OCR = RapidOCR()
    return _SHARED_OCR

def get_embedding_function():
    from langchain_huggingface import HuggingFaceEmbeddings
    return HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL_NAME,
        model_kwargs={"device": EMBEDDING_DEVICE},
        encode_kwargs={"normalize_embeddings": True, "batch_size": 64},
    )

def _process_pdf_page(args: tuple) -> tuple:
    page_idx, pdf_bytes, filename, enable_ocr = args
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc[page_idx]
    text = page.get_text().strip()
    
    if enable_ocr and (not text or len(text) < 20):
        ocr = get_ocr()
        pix = page.get_pixmap(dpi=80, colorspace=fitz.csGRAY)
        img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width)
        ocr_results, _ = ocr(img_array)
        if ocr_results:
            text = " ".join([line[1] for line in ocr_results]).strip()
            
    doc.close()
    return page_idx, text, filename

def extract_documents(
    uploaded_files: List[Any],
    start_page: int = 1,
    end_page: Optional[int] = None,
    enable_ocr: bool = True,
    progress_callback: Optional[Callable[[str, float], None]] = None,
) -> List[Document]:
    documents: List[Document] = []
    total_files = len(uploaded_files)

    for f_idx, file in enumerate(uploaded_files):
        filename = file.name
        ext = Path(filename).suffix.lower()
        file.seek(0)
        file_bytes = file.read()

        if ext == ".pdf":
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            total_pages = len(doc)
            doc.close()

            actual_start = max(0, start_page - 1)
            actual_end = min(total_pages, end_page) if end_page else total_pages
            target_pages = list(range(actual_start, actual_end))
            pages_to_process = len(target_pages)

            if enable_ocr:
                get_ocr()

            tasks = [(idx, file_bytes, filename, enable_ocr) for idx in target_pages]
            results = []
            completed = 0

            with ThreadPoolExecutor(max_workers=min(6, os.cpu_count() or 4)) as executor:
                futures = {executor.submit(_process_pdf_page, t): t for t in tasks}
                for fut in as_completed(futures):
                    page_idx, text, fname = fut.result()
                    completed += 1
                    if progress_callback:
                        progress = (f_idx + completed / max(1, pages_to_process)) / total_files * 0.8
                        progress_callback(f"Reading PDF '{fname}' (Page {page_idx + 1})...", progress)
                    if text:
                        results.append((page_idx, text, fname))

            results.sort(key=lambda x: x[0])
            for p_idx, text, fname in results:
                documents.append(Document(page_content=text, metadata={"source": fname, "page": p_idx + 1}))

        elif ext == ".docx":
            import docx
            doc = docx.Document(BytesIO(file_bytes))
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if full_text:
                documents.append(Document(page_content=full_text, metadata={"source": filename, "page": 1}))

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            for s_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    documents.append(Document(page_content="\n".join(slide_text), metadata={"source": filename, "page": s_idx}))

        elif ext in [".csv", ".xlsx"]:
            if ext == ".csv":
                df = pd.read_csv(BytesIO(file_bytes))
            else:
                df = pd.read_excel(BytesIO(file_bytes))
            csv_text = df.to_string(index=False)
            documents.append(Document(page_content=csv_text, metadata={"source": filename, "page": 1}))

        else:
            try:
                content = file_bytes.decode("utf-8", errors="ignore").strip()
                if content:
                    documents.append(Document(page_content=content, metadata={"source": filename, "page": 1}))
            except Exception:
                pass

    return documents

def chunk_documents(documents: List[Document]) -> List[Document]:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", " ", ""],
    )
    return splitter.split_documents(documents)

def get_user_db_path(user_id: str = "default") -> str:
    safe_user = "".join([c if c.isalnum() else "_" for c in str(user_id)]).lower()
    path = os.path.join(BASE_USER_DATA_DIR, safe_user, "chroma_db")
    os.makedirs(path, exist_ok=True)
    return path

def create_session_vector_store(session_id: str, chunks: List[Document], embeddings: Any, user_id: str = "default", **kwargs):
    from langchain_chroma import Chroma
    db_path = get_user_db_path(user_id)
    collection_name = f"col_{session_id.replace('-', '_')}"
    return Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        collection_name=collection_name,
        persist_directory=db_path,
    )

def load_session_vector_store(session_id: str, embeddings: Any, user_id: str = "default", **kwargs):
    db_path = get_user_db_path(user_id)
    if os.path.exists(db_path) and os.listdir(db_path):
        try:
            from langchain_chroma import Chroma
            collection_name = f"col_{session_id.replace('-', '_')}"
            store = Chroma(
                collection_name=collection_name,
                embedding_function=embeddings,
                persist_directory=db_path,
            )
            if store._collection.count() > 0:
                return store
        except Exception:
            return None
    return None

def delete_session_vector_store(session_id: str, user_id: str = "default", embeddings: Any = None, **kwargs):
    try:
        from langchain_chroma import Chroma
        db_path = get_user_db_path(user_id)
        collection_name = f"col_{session_id.replace('-', '_')}"
        emb = embeddings or get_embedding_function()
        store = Chroma(
            collection_name=collection_name,
            embedding_function=emb,
            persist_directory=db_path,
        )
        store.delete_collection()
    except Exception:
        pass

class LocalHybridRetriever:
    def __init__(self, vector_store: Any, chunks: Optional[List[Document]] = None, k: int = 10):
        self.vector_store = vector_store
        self.k = k
        self.bm25 = None
        if chunks:
            try:
                from langchain_community.retrievers import BM25Retriever
                self.bm25 = BM25Retriever.from_documents(chunks)
                self.bm25.k = k
            except Exception:
                self.bm25 = None

    def invoke(self, query: str) -> List[Document]:
        candidates: List[Document] = []
        seen = set()

        if self.bm25:
            try:
                bm25_docs = self.bm25.invoke(query)
                for d in bm25_docs:
                    if d.page_content not in seen:
                        candidates.append(d)
                        seen.add(d.page_content)
            except Exception:
                pass

        if self.vector_store:
            try:
                dense_docs = self.vector_store.similarity_search(query, k=self.k)
                for d in dense_docs:
                    if d.page_content not in seen:
                        candidates.append(d)
                        seen.add(d.page_content)
            except Exception:
                pass

        return candidates

def build_hybrid_retriever(chunks: List[Document], vector_store: Any):
    return LocalHybridRetriever(vector_store, chunks, k=RETRIEVAL_CANDIDATES)

def rerank_documents(query: str, docs: List[Document], top_k: int = TOP_K) -> List[Document]:
    if not docs:
        return []
    try:
        from flashrank import Ranker, RerankRequest
        ranker = Ranker(model_name="ms-marco-TinyBERT-L-2-v2", cache_dir="./opt_models")
        passages = [{"id": i, "text": doc.page_content, "meta": doc.metadata} for i, doc in enumerate(docs)]
        req = RerankRequest(query=query, passages=passages)
        results = ranker.rerank(req)
        
        reranked_docs = []
        for res in results[:top_k]:
            orig_doc = docs[res["id"]]
            reranked_docs.append(orig_doc)
        return reranked_docs
    except Exception:
        return docs[:top_k]

def contextualize_query(query: str, chat_history: List[Dict[str, Any]]) -> str:
    if not chat_history or len(chat_history) < 2:
        return query
    
    from langchain_ollama import ChatOllama
    from langchain_core.prompts import PromptTemplate
    
    history_snippet = "\n".join([f"{m['role'].capitalize()}: {m['content']}" for m in chat_history[-4:]])
    
    prompt = PromptTemplate.from_template(
        """Given a chat history and the latest user question which might reference context in the chat history, formulate a standalone search question which can be understood without the chat history. Do NOT answer the question, just reformulate it if needed and otherwise return it as is.

Chat History:
{history}

Latest Question: {question}
Standalone Question:"""
    )
    
    llm = ChatOllama(model=OLLAMA_MODEL, temperature=0.0)
    chain = prompt | llm
    try:
        standalone = chain.invoke({"history": history_snippet, "question": query}).content.strip()
        return standalone if standalone else query
    except Exception:
        return query

def format_docs(docs: List[Document]) -> str:
    return "\n\n".join(
        f"[Passage {i} | Source: {doc.metadata.get('source', 'Unknown')} (Page {doc.metadata.get('page', 'N/A')})]\n{doc.page_content.strip()}"
        for i, doc in enumerate(docs, 1)
    )

RAG_PROMPT_TEMPLATE = """You are an accurate, helpful AI research assistant.
Answer the question using strictly the retrieved context below.
If the context does not contain the answer, say "I cannot find the answer based on the provided documents."

Context:
{context}

Question:
{question}

Answer:"""

def generate_study_tool(tool_type: str, chunks: List[Document], extra_param: int = 5) -> str:
    from langchain_ollama import ChatOllama

    sample_size = min(15, len(chunks))
    step = max(1, len(chunks) // sample_size)
    sampled_chunks = [chunks[i] for i in range(0, len(chunks), step)][:sample_size]
    context = "\n\n".join([f"[Source: {d.metadata.get('source')} (Page {d.metadata.get('page')})]\n{d.page_content}" for d in sampled_chunks])

    if tool_type == "quiz":
        prompt_text = f"""Based on the provided document content, generate an interactive Multiple Choice Quiz consisting of {extra_param} questions.
Format EACH question clearly as follows:
### Question X: [Question text]
- A) [Option A]
- B) [Option B]
- C) [Option C]
- D) [Option D]

**Correct Answer:** [Letter]
**Explanation:** [Brief reason citing the text]

Document Content:
{context}"""

    elif tool_type == "flashcards":
        prompt_text = f"""Based on the provided document content, extract the top {extra_param} most important Technical Terms, Formulas, and Core Definitions into study flashcards.
Format EACH flashcard clearly as:
### 🎴 Flashcard X: [Concept / Term / Formula Name]
- **Definition / Principle:** [Clear 2-sentence explanation]
- **Key Formula / Application:** [Formula or real-world application if applicable]
- **Source Context:** [Page or topic reference]

Document Content:
{context}"""

    elif tool_type == "summary":
        prompt_text = f"""Generate a structured, comprehensive Executive Summary of the provided document content.
Include:
1. **Executive Overview**: High-level purpose and core theme.
2. **Key Subject Modules & Concepts**: Bulleted breakdown of major topics.
3. **Important Technical Principles & Formulas**: Key technical takeaways.
4. **Summary Conclusion**: Final takeaways.

Document Content:
{context}"""

    else:
        return "Invalid study tool requested."

    llm = ChatOllama(model=OLLAMA_MODEL, temperature=0.3)
    response = llm.invoke(prompt_text)
    return response.content
